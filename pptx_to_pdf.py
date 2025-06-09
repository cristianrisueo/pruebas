from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import requests
import io


def pptx_url_to_pdf_bytes(pptx_url):
    """
    Descarga un PPTX desde una URL y lo convierte a PDF en memoria.

    Args:
        pptx_url (str): URL del archivo PPTX

    Returns:
        bytes: El PDF generado como bytes
    """

    # Descargar el PPTX desde la URL
    response = requests.get(pptx_url)
    response.raise_for_status()  # Lanza excepción si hay error HTTP

    # Cargar el PPTX desde los bytes descargados
    pptx_bytes = io.BytesIO(response.content)
    presentation = Presentation(pptx_bytes)

    # Crear buffer en memoria para el PDF
    pdf_buffer = io.BytesIO()

    # Crear el PDF en memoria
    c = canvas.Canvas(pdf_buffer, pagesize=letter)
    width, height = letter

    # Procesar cada diapositiva
    for slide_num, slide in enumerate(presentation.slides):
        if slide_num > 0:
            c.showPage()

        # Título de la diapositiva
        y_position = height - 50
        c.setFont("Helvetica-Bold", 16)
        c.drawString(50, y_position, f"Diapositiva {slide_num + 1}")

        # Extraer texto
        y_position -= 40
        c.setFont("Helvetica", 12)

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    lines = text.split("\n")
                    for line in lines:
                        if y_position < 50:
                            c.showPage()
                            y_position = height - 50

                        c.drawString(50, y_position, line[:80])
                        y_position -= 20

    # Guardar el PDF en el buffer
    c.save()

    # Obtener los bytes del PDF
    pdf_bytes = pdf_buffer.getvalue()
    pdf_buffer.close()

    return pdf_bytes


# Ejemplo de uso (descomenta para probar)
if __name__ == "__main__":
    url_pptx = "https://scholar.harvard.edu/files/torman_personal/files/samplepptx.pptx"
    try:
        pdf_bytes = pptx_url_to_pdf_bytes(url_pptx)
        print(f"✅ PDF generado: {len(pdf_bytes)} bytes")

        # Opcional: guardar los bytes en un archivo
        with open("outputs/pdf/desde_url.pdf", "wb") as f:
            f.write(pdf_bytes)
    except Exception as e:
        print(f"❌ Error al convertir desde URL: {e}")
